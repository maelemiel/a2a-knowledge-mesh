#!/usr/bin/env python3
"""Generate A2A Knowledge Mesh pitch deck (.pptx)."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Modern Slate & Cyan SaaS Palette
BG_DARK = RGBColor(0x0F, 0x17, 0x2A)     # Slate 900
ACCENT = RGBColor(0x06, 0xB6, 0xD4)      # Cyan 500
WHITE = RGBColor(0xF8, 0xFA, 0xFC)      # Slate 50
GRAY = RGBColor(0x94, 0xA3, 0xB8)       # Slate 400
DARK_CARD = RGBColor(0x1E, 0x29, 0x3B)  # Slate 800
GREEN = RGBColor(0x10, 0xB9, 0x81)      # Emerald 500
RED = RGBColor(0xF4, 0x3F, 0x5E)        # Rose 500
AMBER = RGBColor(0xF5, 0x9E, 0x0B)      # Amber 500
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)     # Violet 500

FONT_NAME = "Inter"


def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name=FONT_NAME):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_bullet_frame(slide, left, top, width, height, items, size=14, color=GRAY):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = FONT_NAME
        p.space_after = Pt(6)
    return txBox


def cover(slide):
    add_bg(slide)
    # Accent bar top
    add_shape(slide, Inches(0), Inches(0), Inches(13.33), Pt(4), ACCENT)
    # Tag
    add_text_box(slide, Inches(0.8), Inches(0.6), Inches(6), Inches(0.5),
                 "BAND OF AGENTS · HACKATHON 2026", size=11, color=ACCENT, bold=True)
    # Title
    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(9), Inches(1.2),
                 "A2A Knowledge Mesh", size=44, color=WHITE, bold=True)
    # Subtitle
    add_text_box(slide, Inches(0.8), Inches(2.6), Inches(10), Inches(0.8),
                 "Catch enterprise knowledge drift before it catches you.",
                 size=20, color=GRAY, bold=False)
    
    # Authors
    add_text_box(slide, Inches(0.8), Inches(3.3), Inches(6), Inches(0.4),
                 "Created by: Eliott Raguin & Maël Perrigaud", size=13, color=ACCENT, bold=True)
    
    # Modern Agent Badge Cards
    agents = ["Scraper", "Keeper", "Reconciler", "Registry", "Bridge"]
    colors = [ACCENT, GREEN, AMBER, PURPLE, RED]
    descriptions = ["Data Extraction", "Data Storage", "Conflict Resolution", "Data Indexing", "System Overview"]
    icons = ["🔍", "💾", "⚖️", "📋", "📊"]
    
    for i, (a, c, d, ic) in enumerate(zip(agents, colors, descriptions, icons)):
        x = Inches(0.8 + i * 2.38)
        y = Inches(4.2)
        # Background card
        add_shape(slide, x, y, Inches(2.2), Inches(2.0), DARK_CARD, c)
        
        # Icon
        add_text_box(slide, x, y + Inches(0.2), Inches(2.2), Inches(0.4),
                     ic, size=24, color=WHITE, align=PP_ALIGN.CENTER)
        
        # Name
        add_text_box(slide, x, y + Inches(0.8), Inches(2.2), Inches(0.3),
                     a, size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
                     
        # Short description
        add_text_box(slide, x + Inches(0.1), y + Inches(1.25), Inches(2.0), Inches(0.6),
                     d, size=11, color=GRAY, align=PP_ALIGN.CENTER)

    # Bottom
    add_text_box(slide, Inches(0.8), Inches(6.8), Inches(8), Inches(0.4),
                 "Five specialized agents · One Band room · Real-time drift detection",
                 size=12, color=GRAY)


def problem(slide):
    add_bg(slide)
    add_shape(slide, Inches(0), Inches(0), Inches(13.33), Pt(4), RED)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.5),
                 "THE CHALLENGE", size=11, color=RED, bold=True)
    add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.8),
                 "Enterprise Knowledge Drift", size=32, color=WHITE, bold=True)
    add_text_box(slide, Inches(0.8), Inches(1.9), Inches(11.7), Inches(0.7),
                 "Docs say one thing. Code says another. CI lives in its own world. Nobody notices until it breaks.",
                 size=16, color=GRAY)

    # Drift scenario boxes
    scenarios = [
        ("README", "Python 3.9", "pip install"),
        ("pyproject.toml", "Python >=3.11", "uv install"),
        ("CI config", "Python 3.12", "—"),
        ("Dockerfile", "Python 3.10", "—"),
        ("Architecture docs", "Firebase auth", "—"),
        ("Code base", "Supabase auth", "—"),
    ]
    for i, (src, ver, pkg) in enumerate(scenarios):
        col = i % 3
        row = i // 3
        x = Inches(0.8 + col * 4.0)
        y = Inches(2.8 + row * 1.5)
        add_shape(slide, x, y, Inches(3.7), Inches(1.2), DARK_CARD, RED if "auth" in ver or "3.9" in ver or "Firebase" in ver else GRAY)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.12), Inches(3.3), Inches(0.3),
                     src, size=12, color=RED if "auth" in ver or "3.9" in ver or "Firebase" in ver else ACCENT, bold=True)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.45), Inches(3.3), Inches(0.35),
                     ver, size=14, color=WHITE, bold=True)
        if pkg != "—":
            add_text_box(slide, x + Inches(0.2), y + Inches(0.8), Inches(3.3), Inches(0.3),
                         pkg, size=11, color=GRAY)

    add_text_box(slide, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.5),
                 "⚠ Same project. Six contradicting sources. Which one is the actual ground truth?",
                 size=14, color=AMBER, bold=True)


def solution(slide):
    add_bg(slide)
    add_shape(slide, Inches(0), Inches(0), Inches(13.33), Pt(4), GREEN)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.5),
                 "THE SOLUTION", size=11, color=GREEN, bold=True)
    add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.8),
                 "Five Agents, One Room", size=32, color=WHITE, bold=True)
    add_text_box(slide, Inches(0.8), Inches(1.9), Inches(11.7), Inches(0.6),
                 "Band-native multi-agent platform that detects, scores, and resolves knowledge drift in real time.",
                 size=16, color=GRAY)

    agents_detail = [
        ("Scraper", "Reads repos, extracts structured facts via LLM", ACCENT, "🔍"),
        ("Keeper", "SQLite fact store with O(n log n) conflict detection via JOIN", GREEN, "💾"),
        ("Reconciler", "LLM-powered conflict scoring + interactive resolution loop", AMBER, "⚖️"),
        ("Registry", "Dynamic service discovery for mesh agents", PURPLE, "📋"),
        ("Bridge", "Real-time dashboard with modern glassmorphism UI", RED, "📊"),
    ]
    for i, (name, desc, color, icon) in enumerate(agents_detail):
        y = Inches(2.7 + i * 0.9)
        add_shape(slide, Inches(0.8), y, Inches(11.7), Inches(0.78), DARK_CARD, color)
        add_text_box(slide, Inches(1.2), y + Inches(0.18), Inches(2.0), Inches(0.4),
                     f"{icon} {name}", size=16, color=color, bold=True)
        add_text_box(slide, Inches(3.4), y + Inches(0.18), Inches(8.8), Inches(0.4),
                     desc, size=14, color=WHITE)


def architecture(slide):
    add_bg(slide)
    # Thin Cyan bar at top
    add_shape(slide, Inches(0), Inches(0), Inches(13.33), Pt(4), ACCENT)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.5),
                 "ARCHITECTURE", size=11, color=ACCENT, bold=True)
    add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.6),
                 "Collaborative Multi-Agent Flow", size=24, color=WHITE, bold=True)

    # Add the beautiful infographic image showing explanations of each agent
    image_path = "/home/mael/mael/Dev/band/a2a-knowledge-mesh/docs/five_agents_workflow.jpg"
    try:
        slide.shapes.add_picture(image_path, Inches(0.8), Inches(1.8), Inches(8.0), Inches(4.5))
    except Exception:
        # Fallback to drawing a simple card if image loading fails
        add_shape(slide, Inches(0.8), Inches(1.8), Inches(8.0), Inches(4.5), DARK_CARD, ACCENT)
        add_text_box(slide, Inches(1.2), Inches(3.6), Inches(7.2), Inches(1.0),
                     "[Workflow Infographic Diagram]", size=16, color=GRAY, align=PP_ALIGN.CENTER)

    # Explanation card on the right
    x_desc = Inches(9.1)
    y_desc = Inches(1.8)
    w_desc = Inches(3.4)
    add_shape(slide, x_desc, y_desc, w_desc, Inches(4.5), DARK_CARD, ACCENT)

    add_text_box(slide, x_desc + Inches(0.2), y_desc + Inches(0.2), w_desc - Inches(0.4), Inches(0.4),
                 "Key Infrastructure Details", size=14, color=ACCENT, bold=True)

    details = [
        "• Band Room Transport: agents communicate using @mentions and WebSockets.",
        "• Decentralized DB: each agent maintains its own SQLite database.",
        "• A2A Protocol: standard JSON-RPC 2.0 endpoints for local agent requests.",
        "• Provenance Tracked: facts are RDF-lite triples containing subject, predicate, object, source_id, and timestamp."
    ]

    add_bullet_frame(slide, x_desc + Inches(0.2), y_desc + Inches(0.7), w_desc - Inches(0.4), Inches(3.6),
                     details, size=11, color=WHITE)

    # Bottom footnote
    add_text_box(slide, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.4),
                 "Secure bearer authentication + HMAC payload signatures protect all inter-agent operations.",
                 size=11, color=GRAY)


def demo_flow(slide):
    add_bg(slide)
    add_shape(slide, Inches(0), Inches(0), Inches(13.33), Pt(4), AMBER)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.5),
                 "DEMONSTRATION", size=11, color=AMBER, bold=True)
    add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.8),
                 "Watch It In Action", size=32, color=WHITE, bold=True)

    steps = [
        ("1", "Scraper scans repo", "8 facts extracted from docs-repo + code-repo", ACCENT),
        ("2", "Keeper detects conflicts", "3 contradictions found (framework, version, database)", GREEN),
        ("3", "Reconciler scores severity", "critical · high · medium — with LLM confidence ratings", AMBER),
        ("4", "Human reviews & resolves", "Interactive loop via @mentions in Band room", RED),
        ("5", "Dashboard updates live", "Timeline, counters, and mesh graph reflect changes", WHITE),
    ]
    for i, (num, title, desc, color) in enumerate(steps):
        y = Inches(2.0 + i * 0.95)
        add_shape(slide, Inches(0.8), y, Inches(11.7), Inches(0.82), DARK_CARD, color)
        
        # Step Number badge
        num_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.1), y + Inches(0.16), Inches(0.5), Inches(0.5))
        num_shape.fill.solid()
        num_shape.fill.fore_color.rgb = color
        num_shape.line.fill.background()
        tf = num_shape.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.font.name = FONT_NAME
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(2)

        add_text_box(slide, Inches(1.8), y + Inches(0.1), Inches(4.0), Inches(0.4),
                     title, size=16, color=color, bold=True)
        add_text_box(slide, Inches(1.8), y + Inches(0.44), Inches(10.0), Inches(0.35),
                     desc, size=12, color=GRAY)

    add_text_box(slide, Inches(0.8), Inches(6.8), Inches(11.7), Inches(0.3),
                 "Run: bash scripts/run_mesh.sh  ·  Offline mode: uv run python scripts/quick_demo.py",
                 size=12, color=GRAY)


def mesh_graph(slide):
    add_bg(slide)
    add_shape(slide, Inches(0), Inches(0), Inches(13.33), Pt(4), PURPLE)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.5),
                 "VISUALIZATION", size=11, color=PURPLE, bold=True)
    add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.8),
                 "Knowledge Mesh Graph", size=32, color=WHITE, bold=True)

    # Legend buttons
    add_shape(slide, Inches(0.8), Inches(2.0), Inches(2.8), Inches(0.5), DARK_CARD, ACCENT)
    add_text_box(slide, Inches(0.8), Inches(2.05), Inches(2.8), Inches(0.4),
                 "🔵 Subject node", size=12, color=ACCENT, align=PP_ALIGN.CENTER)

    add_shape(slide, Inches(3.8), Inches(2.0), Inches(2.8), Inches(0.5), DARK_CARD, GREEN)
    add_text_box(slide, Inches(3.8), Inches(2.05), Inches(2.8), Inches(0.4),
                 "🟢 Valid fact", size=12, color=GREEN, align=PP_ALIGN.CENTER)

    add_shape(slide, Inches(6.8), Inches(2.0), Inches(2.8), Inches(0.5), DARK_CARD, RED)
    add_text_box(slide, Inches(6.8), Inches(2.05), Inches(2.8), Inches(0.4),
                 "🔴 Active conflict", size=12, color=RED, align=PP_ALIGN.CENTER)

    # Simulated graph description
    items = [
        "• mesh graphify generates an interactive Vis.js network graph",
        "• Click any node to inspect: versions, timestamps, source files",
        "• Subjects (blue) connect to facts (green) — conflicts (red) are visually scannable",
        "• Zoom, pan, filter — built for human-in-the-loop debugging",
    ]
    add_bullet_frame(slide, Inches(0.8), Inches(2.8), Inches(11.7), Inches(1.8), items, size=14, color=WHITE)

    # Sample graph visual
    sample_nodes = [
        (Inches(1.5), Inches(5.1), ACCENT, "project-ALLY"),
        (Inches(3.3), Inches(4.7), GREEN, "framework=Next.js"),
        (Inches(3.3), Inches(5.7), RED, "framework=React"),
        (Inches(5.7), Inches(5.1), GREEN, "language=TypeScript"),
        (Inches(7.9), Inches(4.7), GREEN, "version=18.2"),
        (Inches(7.9), Inches(5.7), RED, "version=16.8"),
    ]
    for x, y, c, label in sample_nodes:
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(2.2), Inches(0.75))
        shape.fill.solid()
        shape.fill.fore_color.rgb = DARK_CARD
        shape.line.color.rgb = c
        shape.line.width = Pt(3)
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(12)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.font.name = FONT_NAME
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(8)


def tech_stack(slide):
    add_bg(slide)
    add_shape(slide, Inches(0), Inches(0), Inches(13.33), Pt(4), ACCENT)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.5),
                 "TECH STACK", size=11, color=ACCENT, bold=True)
    add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.8),
                 "Technology Architecture", size=32, color=WHITE, bold=True)

    stack = [
        ("Runtime & Storage", "Python 3.11+ · Starlette + Uvicorn · SQLite WAL mode"),
        ("Agent Coordination", "Band SDK (WebSocket) · A2A Protocol (JSON-RPC 2.0)"),
        ("LLM Inference", "Featherless / OpenAI API · Resilience: 2x retries + fallback"),
        ("Security & Auth", "HMAC body signing · Per-role bearer tokens · master fallback"),
        ("Developer Tools", "Glassmorphic dashboard UI · Interactive Vis.js network graph"),
        ("Testing & Fixtures", "std unittest · Offline fixtures covering 4 drift scenarios"),
    ]
    for i, (cat, desc) in enumerate(stack):
        col = i % 2
        row = i // 2
        x = Inches(0.8 + col * 6.0)
        y = Inches(2.1 + row * 1.55)
        add_shape(slide, x, y, Inches(5.7), Inches(1.3), DARK_CARD, GRAY)
        add_text_box(slide, x + Inches(0.3), y + Inches(0.18), Inches(5.1), Inches(0.35),
                     cat, size=15, color=ACCENT, bold=True)
        add_text_box(slide, x + Inches(0.3), y + Inches(0.55), Inches(5.1), Inches(0.65),
                     desc, size=13, color=WHITE)


def why_band(slide):
    add_bg(slide)
    add_shape(slide, Inches(0), Inches(0), Inches(13.33), Pt(4), GREEN)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.5),
                 "WHY BAND", size=11, color=GREEN, bold=True)
    add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.8),
                 "Why Band of Agents?", size=32, color=WHITE, bold=True)

    reasons = [
        "• Shared room collaboration: 5 agents communicate seamlessly in one room.",
        "• Structured context hand-off: @mentions transfer structured fact payloads.",
        "• Real-time events: WebSockets keep system dashboard synchronized without polling.",
        "• Human-in-the-loop: Reconciler offers AI recommendations, human stays in control.",
        "• Open standard ready: standard LVK/WebSocket and local A2A fallback options.",
        "• Single-command launch: one unified script starts all agents and database services."
    ]
    add_bullet_frame(slide, Inches(0.8), Inches(2.2), Inches(11.7), Inches(4.5), reasons, size=16, color=WHITE)


def business_value(slide):
    add_bg(slide)
    add_shape(slide, Inches(0), Inches(0), Inches(13.33), Pt(4), PURPLE)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.5),
                 "BUSINESS VALUE", size=11, color=PURPLE, bold=True)
    add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.8),
                 "Market Scope & Value Proposition", size=32, color=WHITE, bold=True)

    # TAM / SAM Card
    x_market = Inches(0.8)
    y_market = Inches(2.0)
    w_market = Inches(5.6)
    h_market = Inches(4.5)
    add_shape(slide, x_market, y_market, w_market, h_market, DARK_CARD, PURPLE)
    add_text_box(slide, x_market + Inches(0.3), y_market + Inches(0.2), w_market - Inches(0.6), Inches(0.4),
                 "Market Size & Opportunity", size=16, color=PURPLE, bold=True)
    
    market_details = [
        "• TAM: $12B+ Enterprise Knowledge Management Market.",
        "• SAM: $1.2B DevOps & CI/CD Integrity Assurance sector.",
        "• Revenue Stream 1: SaaS subscriptions per repository & active user.",
        "• Revenue Stream 2: Self-hosted enterprise licenses for high security compliance.",
        "• High Growth: Drift detection is critical for companies utilizing AI and automated pipelines."
    ]
    add_bullet_frame(slide, x_market + Inches(0.3), y_market + Inches(0.8), w_market - Inches(0.6), Inches(3.4),
                     market_details, size=12, color=WHITE)

    # USP Card (Unique Selling Proposition)
    x_usp = Inches(6.8)
    y_usp = Inches(2.0)
    w_usp = Inches(5.7)
    h_usp = Inches(4.5)
    add_shape(slide, x_usp, y_usp, w_usp, h_usp, DARK_CARD, GREEN)
    add_text_box(slide, x_usp + Inches(0.3), y_usp + Inches(0.2), w_usp - Inches(0.6), Inches(0.4),
                 "Unique Selling Proposition (USP)", size=16, color=GREEN, bold=True)

    usp_details = [
        "• Active Reconciliation: Unlike passive sync tools (e.g. Memory Store, YC P26) which copy raw chat logs, we check all files and detect actual code vs doc contradictions.",
        "• Conflict Scoring: Automatically scores drift severity using local LLM algorithms.",
        "• Human-in-the-loop: Allows easy resolution through Band room commands, keeping developers in control.",
        "• Full Provenance: Tracks source_id, timestamp, and version for every fact and conflict resolved."
    ]
    add_bullet_frame(slide, x_usp + Inches(0.3), y_usp + Inches(0.8), w_usp - Inches(0.6), Inches(3.4),
                     usp_details, size=11, color=WHITE)


def cover_end(slide):
    add_bg(slide)
    add_shape(slide, Inches(0), Inches(0), Inches(13.33), Pt(4), ACCENT)
    add_text_box(slide, Inches(0.8), Inches(2.0), Inches(10), Inches(1.0),
                 "A2A Knowledge Mesh", size=44, color=WHITE, bold=True)
    add_text_box(slide, Inches(0.8), Inches(3.2), Inches(10), Inches(0.6),
                 "github.com/maelemiel/a2a-knowledge-mesh", size=18, color=ACCENT, bold=True)
    add_text_box(slide, Inches(0.8), Inches(4.2), Inches(10), Inches(0.6),
                 "Created by: Eliott Raguin & Maël Perrigaud", size=14, color=WHITE)
    add_text_box(slide, Inches(0.8), Inches(4.7), Inches(10), Inches(0.6),
                 "Built for Band of Agents Hackathon 2026", size=14, color=GRAY)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slides_builders = [
        cover,
        problem,
        solution,
        architecture,
        demo_flow,
        mesh_graph,
        tech_stack,
        why_band,
        business_value,
        cover_end,
    ]

    for builder in slides_builders:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        builder(slide)

    out = "/home/mael/mael/Dev/band/a2a-knowledge-mesh/A2A_Knowledge_Mesh_Deck.pptx"
    prs.save(out)
    print(f"✅ Saved: {out}")
    print(f"   Slides: {len(slides_builders)}")


if __name__ == "__main__":
    main()
